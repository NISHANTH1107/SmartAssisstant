import os, json, hashlib, shutil
import numpy as np
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from dotenv import load_dotenv
from datetime import datetime, timedelta
from db import users_col, chats_col, semantic_cache_col, file_summary_col


# File processing
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

# Wikipedia search
try:
    import wikipedia
    WIKIPEDIA_AVAILABLE = True
except ImportError:
    WIKIPEDIA_AVAILABLE = False
#try OCR
try:
    import pytesseract
    from PIL import Image
    import pdf2image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# LangChain & FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
import google.generativeai as genai

load_dotenv()

DATA_DIR = Path("./data")
CACHE_DIR = DATA_DIR / "cache"
CACHE_DIR.mkdir(parents=True, exist_ok=True)

# Configure Gemini
GEMINI_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_KEY:
    raise ValueError("GEMINI_API_KEY not found in environment variables!")

genai.configure(api_key=GEMINI_KEY)

EMBEDDINGS = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

#semantic cache helpers
def embed_query(text: str) -> np.ndarray:
    return np.array(EMBEDDINGS.embed_query(text.lower().strip()))

def cosine_similarity(vec1: np.ndarray, vec2: np.ndarray) -> float:
    if np.linalg.norm(vec1) == 0 or np.linalg.norm(vec2) == 0:
        return 0.0
    return np.dot(vec1, vec2) / (np.linalg.norm(vec1) * np.linalg.norm(vec2))

#semantic cache
def get_semantic_cache(username: str,chat_name: str,query: str,threshold: float = 0.85) -> Optional[str]:

    query_embedding = embed_query(query)

    candidates = semantic_cache_col.find({
        "username": username,
        "chat_name": chat_name
    })

    best_response = None
    best_score = 0.0

    for doc in candidates:
        cached_embedding = np.array(doc.get("embedding", []))
        if cached_embedding.size == 0:
            continue

        score = cosine_similarity(query_embedding, cached_embedding)

        if score > best_score:
            best_score = score
            best_response = doc["response"]

    if best_score >= threshold:
        print(f"[CACHE HIT][SEMANTIC] score={best_score:.2f} query='{query}'")
        return best_response

    print(f"[CACHE MISS][SEMANTIC] query='{query}'")
    return None


def set_semantic_cache(username: str,chat_name: str,query: str,response: str):
    semantic_cache_col.insert_one({
        "username": username,
        "chat_name": chat_name,
        "query": query,
        "embedding": embed_query(query).tolist(),
        "response": response,
        "timestamp": datetime.utcnow()
    })

    print(f"[CACHE STORE][SEMANTIC] user={username}, chat={chat_name}, query='{query}'")

# #file summary cache
def set_file_summary(file_path: Path, summary: str, username: str):
    file_hash = hashlib.md5(file_path.read_bytes()).hexdigest()

    file_summary_col.update_one(
        {"username": username, "file_path": str(file_path)},
        {
            "$set": {
                "summary": summary,
                "hash": file_hash,
                "timestamp": datetime.utcnow()
            }
        },
        upsert=True
    )

def get_file_summary(file_path: Path, username: str) -> Optional[str]:
    file_hash = hashlib.md5(file_path.read_bytes()).hexdigest()

    doc = file_summary_col.find_one(
        {
            "username": username,
            "file_path": str(file_path),
            "hash": file_hash
        }
    )

    if doc:
        print(f"[CACHE HIT][SUMMARY] user={username}, file={file_path.name}")
        return doc["summary"]

    print(f"[CACHE MISS][SUMMARY] user={username}, file={file_path.name}")
    return None

#wikipedia
def search_wikipedia(query: str, sentences: int = 3) -> str:
    if not WIKIPEDIA_AVAILABLE:
        return ""
    
    try:
        # Search for relevant pages
        search_results = wikipedia.search(query, results=3)
        
        if not search_results:
            return ""
        
        # Get summary of the most relevant page
        try:
            summary = wikipedia.summary(search_results[0], sentences=sentences)
            return f"\n\n=== WIKIPEDIA INFORMATION ===\nTopic: {search_results[0]}\n{summary}\n"
        except wikipedia.exceptions.DisambiguationError as e:
            # If disambiguation, try the first option
            try:
                summary = wikipedia.summary(e.options[0], sentences=sentences)
                return f"\n\n=== WIKIPEDIA INFORMATION ===\nTopic: {e.options[0]}\n{summary}\n"
            except:
                return ""
        except wikipedia.exceptions.PageError:
            return ""
    except Exception:
        return ""

#user-based
def get_user_dirs(username: str):
    user_dir = DATA_DIR / "users" / username
    return {
        'upload_dir': user_dir / "uploads",
        'index_dir': user_dir / "indexes",
    }

def init_user_dirs(username: str):
    dirs = get_user_dirs(username)
    for directory in dirs.values():
        directory.mkdir(parents=True, exist_ok=True)
    return dirs

def cleanup_temp_data(username: str):
    dirs = get_user_dirs(username)
    
    # Clean temp uploads
    temp_upload = dirs['upload_dir'] / "temp"
    if temp_upload.exists():
        shutil.rmtree(temp_upload)
    
    # Clean temp index
    temp_index = dirs['index_dir'] / "temp"
    if temp_index.exists():
        shutil.rmtree(temp_index)

def move_temp_to_chat(chat_name: str, username: str):
    dirs = get_user_dirs(username)
    
    # Move uploads
    temp_upload = dirs['upload_dir'] / "temp"
    chat_upload = dirs['upload_dir'] / chat_name
    if temp_upload.exists():
        if chat_upload.exists():
            # Merge files
            for file in temp_upload.iterdir():
                if file.is_file():
                    shutil.copy2(file, chat_upload / file.name)
            shutil.rmtree(temp_upload)
        else:
            temp_upload.rename(chat_upload)
    
    # Move index
    temp_index = dirs['index_dir'] / "temp"
    chat_index = dirs['index_dir'] / chat_name
    if temp_index.exists():
        if chat_index.exists():
            shutil.rmtree(chat_index)
        temp_index.rename(chat_index)

#file extract
def extract_text_from_file(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    
    try:
        if suffix == ".pdf":
            reader = PdfReader(str(file_path))
            text_parts = []
            
            for page in reader.pages:
                txt = page.extract_text()
                if txt and txt.strip():
                    text_parts.append(txt)
            
            if not text_parts and OCR_AVAILABLE:
                try:
                    images = pdf2image.convert_from_path(str(file_path))
                    for img in images:
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            text_parts.append(ocr_text)
                except Exception:
                    pass
            
            return "\n\n".join(text_parts) if text_parts else "[No text extracted from PDF]"
        
        elif suffix == ".docx":
            doc = DocxDocument(str(file_path))
            return "\n\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        
        elif suffix == ".pptx":
            prs = Presentation(str(file_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n\n".join(texts)
        
        elif suffix == ".txt":
            return file_path.read_text(encoding="utf-8", errors="ignore")
        
        elif suffix in [".xls", ".xlsx"]:
            df = pd.read_excel(file_path)
            return df.to_string()
        
        elif suffix in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"] and OCR_AVAILABLE:
            img = Image.open(file_path)
            return pytesseract.image_to_string(img)
        
    except Exception as e:
        return f"[Error extracting {file_path.name}: {str(e)}]"
    
    return ""

#file management
def save_uploaded_files(uploaded_files, chat_name: str, username: str) -> List[str]:
    saved = []
    dirs = init_user_dirs(username)
    chat_upload_dir = dirs['upload_dir'] / chat_name
    chat_upload_dir.mkdir(exist_ok=True)
    
    for file in uploaded_files:
        file_path = chat_upload_dir / file.name
        with open(file_path, "wb") as f:
            f.write(file.getbuffer())
        saved.append(file.name)
    
    return saved

#FAISS
def build_index_for_chat(chat_name: str, file_names: List[str], username: str):
    dirs = init_user_dirs(username)
    chat_upload_dir = dirs['upload_dir'] / chat_name
    
    all_texts = []
    for file_name in file_names:
        file_path = chat_upload_dir / file_name
        if file_path.exists():
            text = extract_text_from_file(file_path)
            all_texts.append(Document(
                page_content=text,
                metadata={"source": file_name}
            ))
    
    if not all_texts:
        return
    
    # Split into chunks
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    
    chunks = []
    for doc in all_texts:
        splits = splitter.split_text(doc.page_content)
        for split in splits:
            chunks.append(Document(
                page_content=split,
                metadata=doc.metadata
            ))
    
    # Build FAISS index
    if chunks:
        vectorstore = FAISS.from_documents(chunks, EMBEDDINGS)
        index_path = dirs['index_dir'] / chat_name
        vectorstore.save_local(str(index_path))

def load_index(chat_name: str, username: str):
    dirs = get_user_dirs(username)
    index_path = dirs['index_dir'] / chat_name
    if index_path.exists():
        try:
            return FAISS.load_local(str(index_path), EMBEDDINGS, allow_dangerous_deserialization=True)
        except:
            return None
    return None

#file summary
def summarize_file(file_name: str, chat_name: str, username: str) -> str:
    dirs = get_user_dirs(username)
    file_path = dirs['upload_dir'] / chat_name / file_name
    
    if not file_path.exists():
        return f"Error: File '{file_name}' not found."
    
    # Check cache first
    cached_summary = get_file_summary(file_path, username)
    if cached_summary:
        return cached_summary
    
    # Extract text from file
    try:
        text_content = extract_text_from_file(file_path)
        
        if not text_content or text_content.strip() == "":
            return f"Error: No content could be extracted from '{file_name}'."
        
        # Truncate if too long (Gemini has token limits)
        max_chars = 30000
        if len(text_content) > max_chars:
            text_content = text_content[:max_chars] + "\n\n[Content truncated for length...]"
        
        # Create summarization prompt
        prompt = f"""You are an expert study assistant. Please provide a comprehensive summary of the following document.

                Document: {file_name}

                Content:
                {text_content}

                Please provide:
                1. A brief overview (2-3 sentences)
                2. Key points and main topics covered
                3. Important concepts or definitions
                4. Any notable data, figures, or conclusions

                Make the summary clear, well-structured, and educational."""

        # Call Gemini for summarization
        model = genai.GenerativeModel('gemini-2.5-flash-lite')
        response = model.generate_content(prompt)
        summary = response.text
        
        # Cache the summary
        set_file_summary(file_path, summary, username)
        print(f"[CACHE STORE][SUMMARY] user={username}, file={file_path.name}")
        
        return summary
        
    except Exception as e:
        return f"Error generating summary for '{file_name}': {str(e)}"

#generate response
def get_ai_response(query: str, chat_name: str, chat_history: List[Dict], username: str) -> str:  
    # Retrieve context from FAISS
    context_parts = []
    vectorstore = load_index(chat_name, username)
    
    if vectorstore:
        docs = vectorstore.similarity_search(query, k=4)
        if docs:
            context_parts.append("=== RELEVANT CONTENT FROM YOUR FILES ===")
            for i, doc in enumerate(docs, 1):
                context_parts.append(f"\n[Source: {doc.metadata.get('source', 'Unknown')}]")
                context_parts.append(doc.page_content)
    
    # Add conversation history
    if len(chat_history) > 1:
        context_parts.append("\n\n=== CONVERSATION HISTORY ===")
        recent = chat_history[-6:-1] 
        for msg in recent:
            role = msg['role'].upper()
            content = msg['content']
            context_parts.append(f"\n[{role}]: {content}")
    
    context = "\n".join(context_parts)
    
    # Check semantic cache first (now truly semantic!)
    semantic_key = hashlib.md5(
    f"{username}:{chat_name}:{query.lower().strip()}".encode()).hexdigest()

    cached_response = get_semantic_cache(
        username,
        chat_name,
        query
    )
    if cached_response:
        return f"{cached_response}\n\n*[Cached response]*"
    
    
    # Search Wikipedia if no files context
    wiki_info = ""
    if not vectorstore or not context_parts:
        wiki_info = search_wikipedia(query, sentences=3)
        if wiki_info:
            context_parts.append(wiki_info)
            context = "\n".join(context_parts)
    
    # Build prompt
    prompt = f"""You are StudyMate, a helpful AI study assistant. Use the context below to answer the user's question accurately and clearly.

            {context}

            USER QUESTION: {query}

            Provide a clear, well-structured answer. If the context doesn't contain the information, say so and provide a general answer based on your knowledge. Always be helpful and educational."""

    # Call Gemini
    try:
        model = genai.GenerativeModel('gemini-2.5-flash-lite')
        response = model.generate_content(prompt)
        result = response.text
        
        # Cache the response with semantic matching
        set_semantic_cache(username, chat_name, query, result)
        print(f"[CACHE STORE][SEMANTIC] user={username}, chat={chat_name}, query='{query}'")
        
        return result
    except Exception as e:
        return f"I encountered an error while processing your request. Please try again."

#quiz
def generate_quiz_for_chat(chat_name: str, chat_history: List[Dict], username: str, num_questions: int = 10) -> List[Dict]:
    context_parts = []
    vectorstore = load_index(chat_name, username)
    
    if vectorstore:
        all_docs = vectorstore.similarity_search("", k=10)
        for doc in all_docs[:5]:
            context_parts.append(doc.page_content)
    
    for msg in chat_history[-10:]:
        if msg['role'] == 'user':
            context_parts.append(f"Q: {msg['content']}")
        elif msg['role'] == 'assistant' and 'quiz' not in msg:
            context_parts.append(f"A: {msg['content']}")
    
    context = "\n\n".join(context_parts)
    
    prompt = f"""Based on the following study material and conversation, generate a quiz with {num_questions} questions total.

            MATERIAL:
            {context}

            Create a mix of questions:
            - 7-8 Multiple Choice Questions (MCQs)
            - 2-3 Short Answer Questions (2 marks each)

            Use this EXACT JSON format:
            [
            {{
                "type": "mcq",
                "question": "Question text here?",
                "choices": ["Option A", "Option B", "Option C", "Option D"],
                "answer": "B",
                "marks": 1
            }},
            {{
                "type": "short",
                "question": "Define/Explain [concept]?",
                "answer": "A brief 4-5 line answer explaining the concept clearly and concisely.",
                "marks": 2
            }}
            ]

            Guidelines:
            - MCQs should test understanding, application, and recall
            - Short answer questions should ask to "Define", "Explain", "Describe", or "State" key concepts
            - Short answers should be a brief of 4-5 liners at maximum (70-100 words)
            - Focus on important concepts from the material
            - Ensure MCQ answer letters (A/B/C/D) match choice positions
            - Total questions: {num_questions}
            - Maintain a 7:3 or 8:2 ratio (MCQs:Short answers)
            - Interleave question types for variety"""

    try:
        model = genai.GenerativeModel('gemini-2.5-flash-lite')
        response = model.generate_content(prompt)
        
        response_text = response.text.strip()
        
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0].strip()
        
        quiz = json.loads(response_text)
        
        # Validate and ensure we have the right mix
        mcq_count = sum(1 for q in quiz if q.get('type') == 'mcq')
        short_count = sum(1 for q in quiz if q.get('type') == 'short')
        
        # If mix is wrong, filter and adjust
        if mcq_count < 6 or short_count < 1:
            mcqs = [q for q in quiz if q.get('type') == 'mcq'][:mcq_count]
            shorts = [q for q in quiz if q.get('type') == 'short'][:short_count]
            quiz = mcqs + shorts
        
        return quiz[:num_questions]
        
    except Exception as e:
        return [{
            "type": "mcq",
            "question": "Unable to generate quiz at this time. Please try again.",
            "choices": ["Try again later"],
            "answer": "A",
            "marks": 1
        }]

#save chats
def list_chats(username: str) -> List[str]:
    return sorted([
        c["chat_name"]
        for c in chats_col.find(
            {"username": username},
            {"chat_name": 1, "_id": 0}
        )
    ])

def save_chat(chat_name: str, messages: List[Dict], files: List[str], username: str):
    move_temp_to_chat(chat_name, username)

    chats_col.update_one(
        {"username": username, "chat_name": chat_name},
        {
            "$set": {
                "messages": messages,
                "files": files,
                "updated_at": datetime.utcnow()
            }
        },
        upsert=True
    )

def load_chat(chat_name: str, username: str) -> Dict:
    chat = chats_col.find_one(
        {"username": username, "chat_name": chat_name},
        {"_id": 0}
    )
    return chat or {"messages": [], "files": []}

def delete_chat(chat_name: str, username: str):
    chats_col.delete_one(
        {"username": username, "chat_name": chat_name}
    )

    # still delete files + FAISS locally
    dirs = get_user_dirs(username)

    upload_path = dirs['upload_dir'] / chat_name
    if upload_path.exists():
        shutil.rmtree(upload_path)

    index_path = dirs['index_dir'] / chat_name
    if index_path.exists():
        shutil.rmtree(index_path)

def remove_file_from_chat(chat_name: str, file_name: str, username: str):
    chats_col.update_one(
        {"username": username, "chat_name": chat_name},
        {"$pull": {"files": file_name}}
    )

    file_path = get_user_dirs(username)['upload_dir'] / chat_name / file_name
    if file_path.exists():
        file_path.unlink()